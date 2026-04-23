"""Slide extraction and summarization utilities.

This module focuses only on extracting text from PDF/PPTX files and producing
concise summaries. Question-generation has been removed — the app now uses
`ai_slides.summarize_from_file` / `ai_slides.summarize_from_url` exclusively.

Optional dependencies: ``python-pptx``, ``fitz`` (PyMuPDF), and ``openai``.
If unavailable the module falls back to lightweight local heuristics.
"""
from __future__ import annotations

import os
import re
import tempfile
import shutil
from typing import List
from urllib.parse import urlparse

import requests

OPTIONAL_pdf = True
OPTIONAL_pptx = True
OPTIONAL_openai = True
OPTIONAL_transformers = False
try:
    import fitz  # PyMuPDF
except Exception:  # pragma: no cover - optional
    OPTIONAL_pdf = False

try:
    from pptx import Presentation
except Exception:  # pragma: no cover - optional
    OPTIONAL_pptx = False

try:
    import openai
except Exception:  # pragma: no cover - optional
    OPTIONAL_openai = False

# transformers optional and only enabled when user opts in
try:
    import transformers  # type: ignore
    OPTIONAL_transformers = os.environ.get("AI_USE_TRANSFORMERS", "0") == "1"
except Exception:
    OPTIONAL_transformers = False


def _download_to_temp(url: str) -> str:
    parsed = urlparse(url)
    if not parsed.scheme:
        raise ValueError("URL must include scheme (http/https)")

    r = requests.get(url, stream=True, timeout=30)
    r.raise_for_status()

    ext = os.path.splitext(parsed.path)[1] or ""
    fd, tmp = tempfile.mkstemp(suffix=ext)
    os.close(fd)
    with open(tmp, "wb") as fh:
        shutil.copyfileobj(r.raw, fh)
    return tmp


def _extract_text_from_pdf(path: str) -> str:
    if not OPTIONAL_pdf:
        raise RuntimeError("PyMuPDF (fitz) is not installed")
    doc = fitz.open(path)
    parts: List[str] = []
    for page in doc:
        parts.append(page.get_text("text"))
    return "\n\n".join(parts).strip()


def _extract_text_from_pptx(path: str) -> str:
    if not OPTIONAL_pptx:
        raise RuntimeError("python-pptx is not installed")
    prs = Presentation(path)
    parts: List[str] = []
    for slide in prs.slides:
        chunks: List[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                txt = shape.text.strip()
                if txt:
                    chunks.append(txt)
        if chunks:
            parts.append("\n".join(chunks))
    return "\n\n".join(parts).strip()


def extract_text_from_file(path: str) -> str:
    lower = path.lower()
    if lower.endswith(".pdf"):
        return _extract_text_from_pdf(path)
    if lower.endswith(".ppt") or lower.endswith(".pptx"):
        return _extract_text_from_pptx(path)
    # fallback try both
    try:
        return _extract_text_from_pdf(path)
    except Exception:
        try:
            return _extract_text_from_pptx(path)
        except Exception:
            raise RuntimeError("Unsupported file format or missing parsers")


def _local_summarize(text: str, max_chars: int = 2000) -> str:
    s = re.sub(r"\s+", " ", text.strip())
    if len(s) <= max_chars:
        return s
    return s[:max_chars].rsplit(" ", 1)[0] + "..."


def summarize_text(text: str, max_tokens: int = 300) -> str:
    """Summarize extracted text.

    Behavior:
    - If transformers are enabled (opt-in) use a local summarization pipeline.
    - Else if OpenAI API key present, use the API to summarize chunks then merge.
    - Otherwise, return a heuristic paragraph-per-some-slides summary.
    """
    def _chunk_text(s: str, max_chars: int = 40000):
        if not s:
            return []
        parts = [p.strip() for p in re.split(r"\n{2,}", s) if p.strip()]
        if not parts:
            # long blob
            return [s[i : i + max_chars] for i in range(0, len(s), max_chars)]
        chunks: List[str] = []
        cur: List[str] = []
        cur_len = 0
        for p in parts:
            plen = len(p)
            if cur_len + plen + 2 > max_chars and cur:
                chunks.append("\n\n".join(cur))
                cur = [p]
                cur_len = plen
            else:
                cur.append(p)
                cur_len += plen + 2
        if cur:
            chunks.append("\n\n".join(cur))
        return chunks

    # transformers path (opt-in)
    if OPTIONAL_transformers:
        try:
            from transformers import pipeline

            summarizer = pipeline("summarization", model="sshleifer/distilbart-cnn-12-6")
            input_text = text[:200_000]
            chunks = [input_text[i : i + 60_000] for i in range(0, len(input_text), 60_000)]
            parts: List[str] = []
            for chunk in chunks:
                out = summarizer(chunk, max_length=max_tokens, min_length=40, do_sample=False)
                parts.append(out[0]["summary_text"] if isinstance(out, list) and out else str(out))
            return "\n\n".join(parts).strip()
        except Exception:
            pass

    # OpenAI path
    if OPTIONAL_openai and os.environ.get("OPENAI_API_KEY"):
        try:
            openai.api_key = os.environ["OPENAI_API_KEY"]
            chunks = _chunk_text(text, max_chars=80_000)
            chunk_summaries: List[str] = []
            for c in chunks:
                prompt = (
                    "Summarize the following slide content in a concise study-friendly way. "
                    "List key points for students:\n\n" + c
                )
                try:
                    resp = openai.ChatCompletion.create(
                        model="gpt-4o-mini" if hasattr(openai, "ChatCompletion") else "gpt-4o-mini",
                        messages=[{"role": "user", "content": prompt}],
                        max_tokens=max_tokens,
                    )
                    content = resp.choices[0].message.content if hasattr(resp.choices[0], "message") else resp.choices[0].text
                    chunk_summaries.append(content.strip())
                except Exception:
                    chunk_summaries.append(_local_summarize(c, max_chars=1200))

            if len(chunk_summaries) > 1:
                merged = "\n\n".join(chunk_summaries)
                prompt2 = (
                    "Combine and condense the following chunk summaries into a single concise study-friendly summary:\n\n" + merged
                )
                try:
                    resp2 = openai.ChatCompletion.create(
                        model="gpt-4o-mini" if hasattr(openai, "ChatCompletion") else "gpt-4o-mini",
                        messages=[{"role": "user", "content": prompt2}],
                        max_tokens=max_tokens,
                    )
                    content2 = resp2.choices[0].message.content if hasattr(resp2.choices[0], "message") else resp2.choices[0].text
                    return content2.strip()
                except Exception:
                    return _local_summarize(merged, max_chars=2000)

            return chunk_summaries[0] if chunk_summaries else _local_summarize(text)
        except Exception:
            return _local_summarize(text)

    # Heuristic fallback: richer extractive summarization (local, free)
    slides = [s.strip() for s in re.split(r"\n{2,}", text) if s.strip()]
    import html

    # Build sentence list (keep original order and slide mapping)
    sentences: List[str] = []
    slide_map: List[int] = []
    for si, s in enumerate(slides, start=1):
        parts = [p.strip() for p in re.split(r'(?<=[.!?])\s+', s) if p.strip()]
        for p in parts:
            p_clean = html.unescape(p)
            p_clean = re.sub(r"\s+", " ", p_clean).strip()
            if p_clean:
                sentences.append(p_clean)
                slide_map.append(si)

    if not sentences:
        return _local_summarize(text)

    # Simple term-frequency scoring (exclude stopwords)
    stopwords = {
        'the','and','is','in','to','of','a','for','that','on','with','as','are','be','this','by','an','or','from','it','at'
    }
    freq = {}
    for s in sentences:
        for w in re.findall(r"\w+", s.lower()):
            if w in stopwords or len(w) < 3:
                continue
            freq[w] = freq.get(w, 0) + 1

    sent_scores = []
    for s in sentences:
        score = 0
        for w in re.findall(r"\w+", s.lower()):
            score += freq.get(w, 0)
        sent_scores.append(score)

    # choose number of sentences based on requested depth (max_tokens hint)
    target_sentences = min(max(6, int(max_tokens / 30)), len(sentences))

    # pick top-scoring sentences but preserve original order
    ranked_idx = sorted(range(len(sentences)), key=lambda i: sent_scores[i], reverse=True)
    selected_idx = sorted(ranked_idx[:target_sentences])

    # group into paragraphs of 3 sentences each and prefix with slide numbers
    paras: List[str] = []
    per_para = 3
    for i in range(0, len(selected_idx), per_para):
        group_idx = selected_idx[i : i + per_para]
        group_sents = [sentences[j] for j in group_idx]
        slide_nums = sorted({slide_map[j] for j in group_idx})
        slide_label = f"Slides {slide_nums[0]}" if len(slide_nums) == 1 else f"Slides {slide_nums[0]}-{slide_nums[-1]}"
        paras.append(f"{slide_label}: " + " ".join(group_sents))

    # Limit total paragraphs to keep output readable but more in-depth than before
    paras = paras[:8]
    paras = [re.sub(r"\s+", " ", p).strip() for p in paras]

    return "\n\n".join(paras)


def summarize_from_file(path: str, max_tokens: int = 300) -> str:
    text = extract_text_from_file(path)
    return summarize_text(text, max_tokens=max_tokens)


def summarize_from_url(url: str, max_tokens: int = 300) -> str:
    tmp = _download_to_temp(url)
    try:
        text = extract_text_from_file(tmp)
    finally:
        try:
            os.remove(tmp)
        except Exception:
            pass
    return summarize_text(text, max_tokens=max_tokens)


def simplify_text(text: str, max_sentences: int = 8) -> str:
    """Extract short assignment instructions or actionable items from text."""
    if OPTIONAL_openai and os.environ.get("OPENAI_API_KEY"):
        try:
            openai.api_key = os.environ["OPENAI_API_KEY"]
            prompt = (
                "Extract the essential assignment instructions from the following text. "
                "Return a short numbered list of actionable items (due date, deliverables, submission format, minimum requirements).\n\n" + text[:100_000]
            )
            resp = openai.ChatCompletion.create(
                model="gpt-4o-mini" if hasattr(openai, "ChatCompletion") else "gpt-4o-mini",
                messages=[{"role": "user", "content": prompt}],
                max_tokens=400,
            )
            content = resp.choices[0].message.content if hasattr(resp.choices[0], "message") else resp.choices[0].text
            return content.strip()
        except Exception:
            pass

    keywords = [
        "due", "deadline", "submit", "submission", "deliverable", "assignment", "turn in", "required",
        "must", "should", "complete", "due date", "upload", "email", "format",
    ]

    sentences = [s.strip() for s in re.split(r'(?<=[.!?])\s+', text) if s.strip()]
    found: list[str] = []
    for s in sentences:
        low = s.lower()
        if any(k in low for k in keywords):
            if s not in found:
                found.append(s)
        if len(found) >= max_sentences:
            break

    if not found:
        snippet = " ".join(sentences[: max_sentences])
        return _local_summarize(snippet, max_chars=1200)

    out_lines = []
    for i, s in enumerate(found[:max_sentences]):
        out_lines.append(f"{i+1}. {s.strip()}")
    return "\n\n".join(out_lines)


def simplify_from_file(path: str) -> str:
    text = extract_text_from_file(path)
    return simplify_text(text)


def simplify_from_url(url: str) -> str:
    tmp = _download_to_temp(url)
    try:
        text = extract_text_from_file(tmp)
    finally:
        try:
            os.remove(tmp)
        except Exception:
            pass
    return simplify_text(text)
